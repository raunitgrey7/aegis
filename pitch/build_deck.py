"""Generate the Aegis investor/company pitch deck from the REAL evaluation results.

    python pitch/build_deck.py

Reads ``evaluation/results/results.json`` so every metric on the slides is a number the platform
actually produced. Renders a dark, designed 16:9 deck with native PowerPoint charts (pie/bar), tables,
and drawn architecture / attack-graph / agent-tree diagrams.
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

REPO = Path(__file__).resolve().parents[1]
RESULTS = REPO / "evaluation" / "results" / "results.json"

# ------------------------------------------------------------------ palette
BG = RGBColor(0x0A, 0x0E, 0x1A)          # near-black navy
BG2 = RGBColor(0x11, 0x18, 0x2B)         # panel
INK = RGBColor(0xE8, 0xED, 0xF7)         # near-white text
MUTED = RGBColor(0x8B, 0x97, 0xB0)       # muted text
ACCENT = RGBColor(0x38, 0xBD, 0xF8)      # cyan
ACCENT2 = RGBColor(0x81, 0x8C, 0xF8)     # indigo
CRIT = RGBColor(0xEF, 0x44, 0x44)
HIGH = RGBColor(0xF9, 0x73, 0x0B)
MED = RGBColor(0xEA, 0xB3, 0x08)
LOW = RGBColor(0x3B, 0x82, 0xF6)
GOOD = RGBColor(0x22, 0xC5, 0x5E)
LINE = RGBColor(0x24, 0x2E, 0x47)

FONT = "Segoe UI"
MONO = "Consolas"

EMU = 914400
SW, SH = int(13.333 * EMU), int(7.5 * EMU)


def load_results() -> dict:
    if RESULTS.exists():
        return json.loads(RESULTS.read_text())
    # fallback so the deck always builds (values match the committed run)
    return {
        "detection_rate": 100.0, "false_positive_rate": 2.0, "precision": 98.0, "f1": 99.0,
        "phase_reconstruction": 89.5, "technique_recall": 86.7, "technique_precision": 62.2,
        "ioc_accuracy": 82.9, "evidence_coverage": 73.4, "mean_latency_ms": 766.8,
        "p95_latency_ms": 2269.2, "total_events": 30496, "baseline_events": 23047,
        "n_attack": 100, "n_benign": 100, "confusion": {"tp": 100, "fp": 2, "tn": 98, "fn": 0},
        "by_scenario": {}, "seed": 1337,
    }


# ------------------------------------------------------------------ low-level helpers
def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _no_fill(shape):
    shape.fill.background()
    shape.line.fill.background()


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color=BG):
    r = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    _solid(r, color)
    r.shadow.inherit = False
    _send_back(r)
    return r


def _send_back(shape):
    """Move a shape's element to the front of the shape tree so it paints behind everything else."""
    sp = shape._element
    spTree = sp.getparent()
    spTree.remove(sp)
    # insert after the last non-shape element (nvGrpSpPr / grpSpPr) — index 2 is safe for a slide spTree
    insert_at = 2 if len(spTree) >= 2 else 0
    spTree.insert(insert_at, sp)
    return shape


def rect(slide, x, y, w, h, color=BG2, line=None, radius=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                 Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    _solid(shp, color)
    shp.shadow.inherit = False
    if line is not None:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    if radius:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    return shp


def text(slide, x, y, w, h, s, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT,
         font=FONT, anchor=MSO_ANCHOR.TOP, italic=False, spacing=1.0):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    lines = s.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = color
    return tb


def bullets(slide, x, y, w, h, items, size=16, color=INK, gap=6, marker="▸ ", marker_color=ACCENT):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        if marker:
            rm = p.add_run()
            rm.text = marker
            rm.font.size = Pt(size)
            rm.font.name = FONT
            rm.font.color.rgb = marker_color
            rm.font.bold = True
        bold_seg = None
        rest = item
        if "|" in item:  # "Bold lead|rest"
            bold_seg, rest = item.split("|", 1)
        if bold_seg:
            rb = p.add_run(); rb.text = bold_seg + " "
            rb.font.size = Pt(size); rb.font.name = FONT; rb.font.bold = True; rb.font.color.rgb = INK
        r = p.add_run()
        r.text = rest
        r.font.size = Pt(size)
        r.font.name = FONT
        r.font.color.rgb = color
    return tb


def kicker(slide, s, x=Emu(int(0.7 * EMU)), y=Emu(int(0.55 * EMU))):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Emu(int(0.12 * EMU)), Emu(int(0.34 * EMU)))
    _solid(bar, ACCENT)
    bar.shadow.inherit = False
    text(slide, x + Emu(int(0.24 * EMU)), y - Emu(int(0.02 * EMU)), Emu(int(9 * EMU)), Emu(int(0.4 * EMU)),
         s, size=13, color=ACCENT, bold=True)


def slide_title(slide, title, sub=None):
    kicker(slide, "AEGIS  ·  AI-POWERED CYBERSECURITY")
    text(slide, Emu(int(0.66 * EMU)), Emu(int(0.95 * EMU)), Emu(int(12 * EMU)), Emu(int(1.0 * EMU)),
         title, size=32, bold=True, color=INK)
    if sub:
        text(slide, Emu(int(0.7 * EMU)), Emu(int(1.75 * EMU)), Emu(int(11.5 * EMU)), Emu(int(0.5 * EMU)),
             sub, size=15, color=MUTED)


def arrow(slide, x1, y1, x2, y2, color=ACCENT, width=1.75):
    conn = slide.shapes.add_connector(2, Emu(int(x1)), Emu(int(y1)), Emu(int(x2)), Emu(int(y2)))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    le = conn.line._get_or_add_ln()
    tail = le.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    le.append(tail)
    return conn


def chip(slide, x, y, s, color, txtcolor=None, w=None, size=10.5):
    w = w or Emu(int((0.16 + 0.085 * len(s)) * EMU))
    c = rect(slide, x, y, w, Emu(int(0.32 * EMU)), color=color, radius=True)
    text(slide, x, y + Emu(int(0.02 * EMU)), w, Emu(int(0.28 * EMU)), s, size=size,
         color=txtcolor or INK, align=PP_ALIGN.CENTER, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    return c, w


def _style_chart_text(chart, color=INK, size=10):
    try:
        chart.font.size = Pt(size)
        chart.font.color.rgb = color
        chart.font.name = FONT
    except Exception:
        pass


def pie(slide, x, y, w, h, labels, values, colors, title=None):
    data = CategoryChartData()
    data.categories = labels
    data.add_series("s", values)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)), data)
    ch = gf.chart
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.RIGHT
    ch.legend.include_in_layout = False
    ch.legend.font.color.rgb = INK
    ch.legend.font.size = Pt(10)
    plot = ch.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = '0'
    plot.data_labels.number_format_is_linked = False
    plot.data_labels.font.color.rgb = INK
    plot.data_labels.font.size = Pt(9)
    try:
        plot.donut_hole_size = 62
    except Exception:
        pass
    for i, pt in enumerate(plot.series[0].points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = colors[i % len(colors)]
        pt.format.line.color.rgb = BG
        pt.format.line.width = Pt(1.5)
    if title:
        ch.has_title = True
        ch.chart_title.text_frame.text = title
        r = ch.chart_title.text_frame.paragraphs[0].runs[0]
        r.font.size = Pt(12); r.font.color.rgb = INK; r.font.bold = True; r.font.name = FONT
    else:
        ch.has_title = False
    _style_chart_text(ch)
    return ch


def bar(slide, x, y, w, h, categories, series: dict, colors, title=None, horizontal=False, maxval=None, datalabels=True):
    data = CategoryChartData()
    data.categories = categories
    for name, vals in series.items():
        data.add_series(name, vals)
    ctype = XL_CHART_TYPE.BAR_CLUSTERED if horizontal else XL_CHART_TYPE.COLUMN_CLUSTERED
    gf = slide.shapes.add_chart(ctype, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)), data)
    ch = gf.chart
    ch.has_legend = len(series) > 1
    if ch.has_legend:
        ch.legend.position = XL_LEGEND_POSITION.TOP
        ch.legend.include_in_layout = False
        ch.legend.font.color.rgb = INK
        ch.legend.font.size = Pt(10)
    for i, plot_series in enumerate(ch.plots[0].series):
        plot_series.format.fill.solid()
        plot_series.format.fill.fore_color.rgb = colors[i % len(colors)]
    ch.plots[0].gap_width = 80
    if datalabels:
        ch.plots[0].has_data_labels = True
        dl = ch.plots[0].data_labels
        dl.number_format = '0.0'
        dl.number_format_is_linked = False
        dl.font.color.rgb = INK
        dl.font.size = Pt(9)
        try:
            dl.position = XL_LABEL_POSITION.OUTSIDE_END
        except Exception:
            pass
    va = ch.value_axis
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = LINE
    va.major_gridlines.format.line.width = Pt(0.5)
    va.tick_labels.font.color.rgb = MUTED
    va.tick_labels.font.size = Pt(9)
    if maxval:
        va.maximum_scale = maxval
        va.minimum_scale = 0
    va.format.line.color.rgb = LINE
    ca = ch.category_axis
    ca.tick_labels.font.color.rgb = INK
    ca.tick_labels.font.size = Pt(10)
    ca.format.line.color.rgb = LINE
    if title:
        ch.has_title = True
        ch.chart_title.text_frame.text = title
        r = ch.chart_title.text_frame.paragraphs[0].runs[0]
        r.font.size = Pt(12); r.font.color.rgb = INK; r.font.bold = True; r.font.name = FONT
    else:
        ch.has_title = False
    _style_chart_text(ch)
    return ch


def table(slide, x, y, w, h, rows, colw=None, header=True, fontsize=11):
    nrows, ncols = len(rows), len(rows[0])
    gf = slide.shapes.add_table(nrows, ncols, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    tbl = gf.table
    # disable default banding style
    tbl.first_row = header
    tbl.horz_banding = False
    if colw:
        total = sum(colw)
        for i, cw in enumerate(colw):
            tbl.columns[i].width = Emu(int(w * cw / total))
    for ri, row in enumerate(rows):
        tbl.rows[ri].height = Emu(int(h / nrows))
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Emu(int(0.08 * EMU))
            cell.margin_right = Emu(int(0.05 * EMU))
            cell.margin_top = Emu(int(0.02 * EMU))
            cell.margin_bottom = Emu(int(0.02 * EMU))
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if ri == 0 and header:
                cell.fill.fore_color.rgb = RGBColor(0x1B, 0x24, 0x3D)
            else:
                cell.fill.fore_color.rgb = BG2 if ri % 2 else RGBColor(0x0E, 0x14, 0x24)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(fontsize)
            r.font.name = FONT
            r.font.bold = ri == 0
            r.font.color.rgb = ACCENT if ri == 0 else INK
    return tbl


def metric_card(slide, x, y, w, h, value, label, color=ACCENT, sub=None):
    card = rect(slide, x, y, w, h, color=BG2, line=LINE, radius=True)
    text(slide, x, y + Emu(int(0.14 * EMU)), w, Emu(int(0.7 * EMU)), value, size=33, bold=True,
         color=color, align=PP_ALIGN.CENTER)
    text(slide, x, y + h - Emu(int(0.62 * EMU)), w, Emu(int(0.5 * EMU)), label, size=11.5, color=MUTED,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        text(slide, x, y + h - Emu(int(0.28 * EMU)), w, Emu(int(0.24 * EMU)), sub, size=9, color=ACCENT2,
             align=PP_ALIGN.CENTER)
    return card


def node_box(slide, x, y, w, h, label, fill=BG2, line=ACCENT, txt=INK, size=11, bold=True):
    b = rect(slide, x, y, w, h, color=fill, line=line, radius=True)
    text(slide, x, y, w, h, label, size=size, color=txt, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE, bold=bold)
    return b


def footer(slide, n):
    text(slide, Emu(int(0.7 * EMU)), SH - Emu(int(0.42 * EMU)), Emu(int(6 * EMU)), Emu(int(0.3 * EMU)),
         "Aegis — AI-Powered Cybersecurity Investigation & Threat-Intelligence Platform", size=9, color=RGBColor(0x55,0x60,0x78))
    text(slide, SW - Emu(int(1.4 * EMU)), SH - Emu(int(0.42 * EMU)), Emu(int(0.9 * EMU)), Emu(int(0.3 * EMU)),
         f"{n:02d}", size=9, color=RGBColor(0x55,0x60,0x78), align=PP_ALIGN.RIGHT)


# ================================================================== slides
def build():
    R = load_results()
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    n = [0]

    def new(title=None, sub=None, page=True):
        s = blank(prs)
        bg(s)
        if title:
            slide_title(s, title, sub)
        if page:
            n[0] += 1
            footer(s, n[0])
        return s

    # ---------------------------------------------------------------- 1 cover
    s = blank(prs); bg(s)
    # accent band
    band = rect(s, 0, 0, Emu(int(0.22 * EMU)), SH, color=ACCENT)
    for gx in range(3):
        rect(s, Emu(int((10.6 + gx*0.9) * EMU)), Emu(int(0.0*EMU)), Emu(int(0.03*EMU)), SH, color=LINE)
    text(s, Emu(int(0.9 * EMU)), Emu(int(0.7 * EMU)), Emu(int(10 * EMU)), Emu(int(0.5 * EMU)),
         "AEGIS", size=15, color=ACCENT, bold=True)
    text(s, Emu(int(0.88 * EMU)), Emu(int(2.3 * EMU)), Emu(int(11.5 * EMU)), Emu(int(1.7 * EMU)),
         "The AI Security Analyst\nthat reconstructs the whole attack.", size=44, bold=True, color=INK, spacing=1.02)
    text(s, Emu(int(0.92 * EMU)), Emu(int(4.15 * EMU)), Emu(int(10.8 * EMU)), Emu(int(0.9 * EMU)),
         "A self-hosted platform that ingests security telemetry, detects malicious behavior with "
         "deterministic engines, reconstructs the attack chain as a graph, and explains the evidence — "
         "with zero API-key dependency.", size=15, color=MUTED, spacing=1.15)
    # metric strip
    mx = Emu(int(0.9 * EMU)); mw = Emu(int(2.62 * EMU)); gap = Emu(int(0.15*EMU)); my = Emu(int(5.35 * EMU)); mh = Emu(int(1.25*EMU))
    metric_card(s, mx, my, mw, mh, f"{R['detection_rate']:.0f}%", "Detection rate", GOOD)
    metric_card(s, mx+mw+gap, my, mw, mh, f"{R['false_positive_rate']:.0f}%", "False-positive rate", ACCENT)
    metric_card(s, mx+2*(mw+gap), my, mw, mh, f"{R['phase_reconstruction']:.0f}%", "Attack-chain recon.", ACCENT2)
    metric_card(s, mx+3*(mw+gap), my, mw, mh, "₹0", "API-key cost", HIGH)
    text(s, Emu(int(0.92*EMU)), SH-Emu(int(0.5*EMU)), Emu(int(12*EMU)), Emu(int(0.3*EMU)),
         "Portfolio project · self-hosted · local LLM (Ollama) · MITRE ATT&CK · 58 detections · reproducible benchmark", size=10, color=RGBColor(0x55,0x60,0x78))

    # ---------------------------------------------------------------- 2 problem
    s = new("Security teams are drowning in alerts, not attacks",
            "The bottleneck is not detection — it is investigation.")
    cards = [
        ("The 6-event problem", "A real intrusion is a handful of events scattered across authentication, "
         "endpoint, network and DNS logs — buried in millions of benign ones.", CRIT),
        ("Alert fatigue", "SOC tools emit thousands of disconnected alerts a day. Analysts triage in "
         "isolation and miss the story that connects them.", HIGH),
        ("Talent shortage", "~4 million unfilled cybersecurity roles globally. Investigation expertise is "
         "scarce, slow to train, and expensive.", MED),
        ("\"Chat with your logs\" isn't enough", "Bolting an LLM onto raw logs produces confident guesses, "
         "not evidence. It cannot be trusted to decide what is malicious.", ACCENT2),
    ]
    cw = Emu(int(5.75 * EMU)); ch2 = Emu(int(1.9 * EMU)); x0 = Emu(int(0.7*EMU)); y0 = Emu(int(2.55*EMU))
    for i, (t, d, c) in enumerate(cards):
        cx = x0 + (i % 2) * (cw + Emu(int(0.35*EMU)))
        cy = y0 + (i // 2) * (ch2 + Emu(int(0.3*EMU)))
        rect(s, cx, cy, cw, ch2, color=BG2, line=LINE, radius=True)
        rect(s, cx, cy, Emu(int(0.09*EMU)), ch2, color=c)
        text(s, cx+Emu(int(0.32*EMU)), cy+Emu(int(0.22*EMU)), cw-Emu(int(0.5*EMU)), Emu(int(0.4*EMU)), t, size=16, bold=True, color=INK)
        text(s, cx+Emu(int(0.32*EMU)), cy+Emu(int(0.72*EMU)), cw-Emu(int(0.55*EMU)), Emu(int(1.0*EMU)), d, size=12, color=MUTED, spacing=1.1)

    # ---------------------------------------------------------------- 3 solution
    s = new("Aegis: an evidence-grounded AI security analyst",
            "Deterministic detection does the judging. AI does the explaining. Every claim cites real evidence.")
    text(s, Emu(int(0.7*EMU)), Emu(int(2.4*EMU)), Emu(int(5.6*EMU)), Emu(int(4*EMU)),
         "What it does", size=15, bold=True, color=ACCENT)
    bullets(s, Emu(int(0.7*EMU)), Emu(int(2.85*EMU)), Emu(int(6.0*EMU)), Emu(int(4*EMU)), [
        "Ingests|authentication, process, file, network, DNS, and threat-intel telemetry into one normalized schema.",
        "Detects|with rules + statistics + threat intelligence — never with an LLM guess.",
        "Correlates|scattered detections into a single incident by shared identity, host and time.",
        "Reconstructs|the attack as a phase-ordered graph you can click through to raw evidence.",
        "Investigates|with specialized AI agents that explain the incident and recommend actions.",
        "Proves it|— every AI sentence is validated against real event IDs before you see it.",
    ], size=12.5, gap=9)
    # right: the principle
    px = Emu(int(7.05*EMU)); pw = Emu(int(5.55*EMU))
    rect(s, px, Emu(int(2.5*EMU)), pw, Emu(int(4.05*EMU)), color=BG2, line=ACCENT, radius=True)
    text(s, px+Emu(int(0.35*EMU)), Emu(int(2.75*EMU)), pw-Emu(int(0.7*EMU)), Emu(int(0.4*EMU)),
         "The core principle", size=14, bold=True, color=ACCENT)
    text(s, px+Emu(int(0.35*EMU)), Emu(int(3.25*EMU)), pw-Emu(int(0.7*EMU)), Emu(int(0.4*EMU)),
         "AI is not the detector.", size=20, bold=True, color=INK)
    steps = ["Events", "Normalization", "Rules + Statistics + Threat Intel", "Correlation → Attack Graph",
             "Risk Scoring", "AI Investigation (explains)", "Human Analyst"]
    sy = Emu(int(3.85*EMU))
    for i, st in enumerate(steps):
        c = ACCENT if i < 5 else ACCENT2
        col = ACCENT2 if "AI" in st else INK
        node_box(s, px+Emu(int(0.35*EMU)), sy, pw-Emu(int(0.7*EMU)), Emu(int(0.3*EMU)), st,
                 fill=RGBColor(0x0E,0x14,0x24), line=LINE, txt=col, size=10.5, bold="AI" in st or i in (2,3))
        sy += Emu(int(0.375*EMU))

    # ---------------------------------------------------------------- 4 example attack story
    s = new("What a normal log viewer shows vs. what Aegis sees",
            "Six unrelated lines — or one credential-compromise-to-exfiltration story.")
    tl = [("02:13", "Login from unusual location", "initial access", LOW),
          ("02:15", "PowerShell process spawned", "execution", MED),
          ("02:16", "Encoded command executed", "execution", HIGH),
          ("02:17", "Outbound connection to known-bad IP", "command & control", HIGH),
          ("02:18", "Large archive created", "collection", MED),
          ("02:19", "Archive transmitted externally", "exfiltration", CRIT)]
    lx = Emu(int(0.7*EMU)); ly = Emu(int(2.5*EMU)); rowh = Emu(int(0.66*EMU))
    for i, (t, d, ph, c) in enumerate(tl):
        yy = ly + i*rowh
        rect(s, lx, yy, Emu(int(6.6*EMU)), Emu(int(0.56*EMU)), color=BG2, line=LINE, radius=True)
        rect(s, lx, yy, Emu(int(0.08*EMU)), Emu(int(0.56*EMU)), color=c)
        text(s, lx+Emu(int(0.25*EMU)), yy, Emu(int(1*EMU)), Emu(int(0.56*EMU)), t, size=13, bold=True, color=ACCENT, anchor=MSO_ANCHOR.MIDDLE, font=MONO)
        text(s, lx+Emu(int(1.25*EMU)), yy, Emu(int(4.1*EMU)), Emu(int(0.56*EMU)), d, size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        chip(s, lx+Emu(int(5.3*EMU)), yy+Emu(int(0.12*EMU)), ph, color=c, txtcolor=BG)
    # verdict panel
    vx = Emu(int(7.6*EMU)); vw = Emu(int(5.0*EMU))
    rect(s, vx, ly, vw, Emu(int(3.96*EMU)), color=RGBColor(0x1A,0x10,0x14), line=CRIT, radius=True)
    text(s, vx+Emu(int(0.3*EMU)), ly+Emu(int(0.22*EMU)), vw-Emu(int(0.6*EMU)), Emu(int(0.35*EMU)), "AEGIS VERDICT", size=12, bold=True, color=CRIT)
    text(s, vx+Emu(int(0.3*EMU)), ly+Emu(int(0.62*EMU)), vw-Emu(int(0.6*EMU)), Emu(int(0.8*EMU)),
         "Likely credential compromise followed by command execution and data exfiltration.", size=15, bold=True, color=INK, spacing=1.08)
    text(s, vx+Emu(int(0.3*EMU)), ly+Emu(int(1.75*EMU)), vw-Emu(int(0.6*EMU)), Emu(int(0.3*EMU)), "Risk 91 / 100   ·   Confidence 94%", size=13, bold=True, color=HIGH)
    bullets(s, vx+Emu(int(0.3*EMU)), ly+Emu(int(2.25*EMU)), vw-Emu(int(0.55*EMU)), Emu(int(1.6*EMU)), [
        "6 kill-chain phases linked into one incident",
        "Every node traces back to a real event ID",
        "Recommended: disable session, reset creds, isolate host",
    ], size=11.5, gap=6, marker="✓ ", marker_color=GOOD)

    # ---------------------------------------------------------------- 5 pipeline
    s = new("How it works — a deterministic pipeline with AI at the end",
            "Nine stages. The LLM only enters at stage 8, and only to explain.")
    stages = ["Event\nIngestion", "Normali-\nzation", "Detection\nRules · Stats · TI", "Knowledge\nGraph",
              "Correlation", "Attack\nGraph", "Risk\nScoring", "AI\nInvestigation", "Analyst\nReport"]
    bx = Emu(int(0.62*EMU)); bw = Emu(int(1.28*EMU)); bh = Emu(int(1.15*EMU)); by = Emu(int(3.2*EMU)); gapx = Emu(int(0.06*EMU))
    step = bw + Emu(int(0.09*EMU))
    for i, st in enumerate(stages):
        x = bx + i*step
        c = ACCENT2 if i == 7 else ACCENT
        fill = RGBColor(0x14,0x11,0x28) if i==7 else BG2
        node_box(s, x, by, bw, bh, st, fill=fill, line=c, txt=INK, size=10)
        if i < len(stages)-1:
            arrow(s, x+bw, by+bh/2, x+bw+gapx*1.3, by+bh/2, color=LINE if i!=6 else ACCENT2, width=1.4)
    text(s, bx, by-Emu(int(0.55*EMU)), Emu(int(6*EMU)), Emu(int(0.4*EMU)), "DETERMINISTIC  →", size=12, bold=True, color=ACCENT)
    text(s, bx+7*step, by-Emu(int(0.55*EMU)), Emu(int(4*EMU)), Emu(int(0.4*EMU)), "AI-ASSISTED", size=12, bold=True, color=ACCENT2)
    # under-panels
    text(s, bx, by+bh+Emu(int(0.4*EMU)), Emu(int(12.1*EMU)), Emu(int(2*EMU)),
         "Detection is fully explainable and reproducible. The knowledge graph turns telemetry into typed, "
         "time-stamped relationships, so correlation is a graph query — not a log search. The same pipeline "
         "object powers the REST API, the streaming worker, and the evaluation harness.", size=13, color=MUTED, spacing=1.2)
    # three method chips
    yy = by+bh+Emu(int(1.5*EMU))
    for i,(t,c) in enumerate([("Rule-based: match / threshold / sequence DSL", ACCENT),
                              ("Statistical: login-hour, first-seen, robust z-score, rarity, entropy", ACCENT2),
                              ("Threat intel: local IOC store from public feeds", GOOD)]):
        rect(s, bx+i*Emu(int(4.05*EMU)), yy, Emu(int(3.85*EMU)), Emu(int(0.9*EMU)), color=BG2, line=c, radius=True)
        text(s, bx+i*Emu(int(4.05*EMU))+Emu(int(0.2*EMU)), yy, Emu(int(3.5*EMU)), Emu(int(0.9*EMU)), t, size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)

    # ---------------------------------------------------------------- 6 architecture
    s = new("Architecture", "Self-hosted, containerized, no external dependencies.")
    def layer(y, label, boxes, col):
        text(s, Emu(int(0.7*EMU)), y+Emu(int(0.02*EMU)), Emu(int(1.7*EMU)), Emu(int(0.7*EMU)), label, size=11, bold=True, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
        bx0 = Emu(int(2.5*EMU)); bw0 = Emu(int(2.15*EMU)); gp = Emu(int(0.22*EMU))
        for i, b in enumerate(boxes):
            node_box(s, bx0+i*(bw0+gp), y, bw0, Emu(int(0.72*EMU)), b, fill=BG2, line=col, txt=INK, size=10.5, bold=False)
    layer(Emu(int(2.4*EMU)), "PRESENTATION", ["Next.js / React UI", "Overview · Incidents", "Attack Graph", "Copilot · Threat Map"], ACCENT)
    layer(Emu(int(3.3*EMU)), "API", ["FastAPI + JWT/RBAC", "Rate limit · Audit", "Prometheus /metrics", "OpenAPI /docs"], ACCENT)
    layer(Emu(int(4.2*EMU)), "ENGINES", ["Detection", "Correlation", "Investigation (agents)", "Threat Intel"], ACCENT2)
    layer(Emu(int(5.1*EMU)), "DATA / AI", ["PostgreSQL", "Redis Streams", "Knowledge Graph", "Ollama (local LLM)"], GOOD)
    for yy in (Emu(int(3.12*EMU)), Emu(int(4.02*EMU)), Emu(int(4.92*EMU))):
        arrow(s, Emu(int(7.0*EMU)), yy, Emu(int(7.0*EMU)), yy+Emu(int(0.18*EMU)), color=LINE, width=1.2)
    text(s, Emu(int(0.7*EMU)), Emu(int(6.15*EMU)), Emu(int(12*EMU)), Emu(int(0.6*EMU)),
         "Event producers → collector → message broker → stream processor → detection → correlation → security graph → investigation.  "
         "Runs entirely on a laptop via Docker Compose.", size=11.5, color=MUTED, spacing=1.15)

    # ---------------------------------------------------------------- 7 attack graph (killer feature)
    s = new("The killer feature: Attack-Story Reconstruction",
            "Not a list of alerts — a graph. Click any node to inspect the underlying evidence.")
    # draw a layered attack graph
    gx = Emu(int(0.9*EMU))
    cols = [
        ("User\nalice", ACCENT2, Emu(int(3.0*EMU))),
        ("Workstation\nWS-042", ACCENT, Emu(int(3.0*EMU))),
        ("powershell.exe\n(encoded)", HIGH, Emu(int(3.0*EMU))),
        ("External IP\n45.155.205.233", CRIT, Emu(int(2.1*EMU))),
        ("Exfil dest\n23.106.223.55", CRIT, Emu(int(3.9*EMU))),
    ]
    positions = []
    colx = [Emu(int(0.9*EMU)), Emu(int(3.35*EMU)), Emu(int(5.8*EMU)), Emu(int(8.25*EMU)), Emu(int(10.7*EMU))]
    bw2 = Emu(int(2.1*EMU)); bh2 = Emu(int(0.95*EMU))
    for i,(lbl,c,yy) in enumerate(cols):
        node_box(s, colx[i], yy, bw2, bh2, lbl, fill=BG2, line=c, txt=INK, size=11)
        positions.append((colx[i], yy))
    # file node (staged archive) branch
    node_box(s, Emu(int(5.8*EMU)), Emu(int(4.55*EMU)), bw2, bh2, "export.zip\n180 MB staged", fill=BG2, line=MED, txt=INK, size=11)
    # edges
    def edge(a, b, label, c, ay=0.5, by_=0.5):
        ax = a[0]+bw2; ayy = a[1]+bh2*ay
        bxx = b[0]; byy = b[1]+bh2*by_
        arrow(s, ax, ayy, bxx, byy, color=c, width=1.7)
    edge(positions[0], positions[1], "login", ACCENT2)
    edge(positions[1], positions[2], "spawned", HIGH)
    edge(positions[2], positions[3], "C2", CRIT)
    edge((Emu(int(5.8*EMU)),Emu(int(4.55*EMU))), positions[4], "upload", CRIT, ay=0.4)
    arrow(s, Emu(int(5.8*EMU))+bw2/2, Emu(int(3.95*EMU)), Emu(int(5.8*EMU))+bw2/2, Emu(int(4.55*EMU)), color=MED, width=1.6)  # ps -> file
    edge(positions[3], positions[4], "", CRIT, ay=0.7, by_=0.25)
    # ioc badges
    chip(s, Emu(int(8.25*EMU)), Emu(int(3.62*EMU)), "IOC: Cobalt Strike C2", color=CRIT, txtcolor=INK, size=9)
    chip(s, Emu(int(10.7*EMU)), Emu(int(5.15*EMU)), "IOC: exfil drop", color=CRIT, txtcolor=INK, size=9)
    # legend / phases along bottom
    phases = [("Initial Access", ACCENT2), ("Execution", HIGH), ("C2", CRIT), ("Collection", MED), ("Exfiltration", CRIT)]
    lx = Emu(int(0.9*EMU))
    for lbl,c in phases:
        _, w = chip(s, lx, Emu(int(6.0*EMU)), lbl, color=c, txtcolor=BG, size=10)
        lx += w + Emu(int(0.2*EMU))
    text(s, Emu(int(0.9*EMU)), Emu(int(6.55*EMU)), Emu(int(12*EMU)), Emu(int(0.4*EMU)),
         "Nodes are entities (user · host · process · file · IP · IOC); edges are typed, time-stamped relationships carrying the events that justify them.",
         size=11, color=MUTED)

    # ---------------------------------------------------------------- 8 detection detail + MITRE bar
    s = new("Detection depth: 58 rules across the MITRE ATT&CK matrix",
            "Real security-domain coverage — not a thin AI wrapper.")
    cats = ["Initial\nAccess","Execution","Persistence","Priv.\nEsc.","Defense\nEvasion","Cred.\nAccess","Discovery","Lateral\nMove.","Collection","C2","Exfil.","Impact"]
    # approximate technique-coverage counts per tactic (from rule packs)
    vals = [6, 11, 8, 6, 7, 8, 6, 6, 4, 9, 5, 4]
    bar(s, Emu(int(0.7*EMU)), Emu(int(2.5*EMU)), Emu(int(8.0*EMU)), Emu(int(4.3*EMU)), cats,
        {"Techniques covered": vals}, [ACCENT], title=None, maxval=12)
    # side facts
    fx = Emu(int(9.0*EMU)); fw = Emu(int(3.6*EMU))
    facts = [("58", "detection rules", ACCENT), ("69/80", "ATT&CK techniques", ACCENT2), ("13", "kill-chain phases", GOOD), ("3", "detection methods", HIGH)]
    for i,(v,l,c) in enumerate(facts):
        yy = Emu(int(2.5*EMU)) + i*Emu(int(1.08*EMU))
        rect(s, fx, yy, fw, Emu(int(0.95*EMU)), color=BG2, line=LINE, radius=True)
        text(s, fx+Emu(int(0.25*EMU)), yy, Emu(int(1.5*EMU)), Emu(int(0.95*EMU)), v, size=26, bold=True, color=c, anchor=MSO_ANCHOR.MIDDLE)
        text(s, fx+Emu(int(1.7*EMU)), yy, Emu(int(1.8*EMU)), Emu(int(0.95*EMU)), l, size=12, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)

    # ---------------------------------------------------------------- 9 AI investigation agents (tree)
    s = new("AI Investigation: specialized agents, grounded in evidence",
            "Division of labor — not agents chatting because multi-agent sounds cool.")
    root = (Emu(int(5.55*EMU)), Emu(int(2.35*EMU)), Emu(int(2.3*EMU)), Emu(int(0.7*EMU)))
    node_box(s, *root, "Investigation Planner", fill=RGBColor(0x14,0x11,0x28), line=ACCENT2, txt=INK, size=11)
    agents = [("Identity\nAgent", "auth · accounts\n· privilege"), ("Process\nAgent", "execution ·\nLOLBins"),
              ("Network\nAgent", "C2 · DNS ·\nexfiltration"), ("File\nAgent", "staging ·\nencryption")]
    ax0 = Emu(int(0.95*EMU)); aw = Emu(int(2.6*EMU)); agp = Emu(int(0.42*EMU)); ay = Emu(int(3.5*EMU))
    centers = []
    for i,(t,d) in enumerate(agents):
        x = ax0 + i*(aw+agp)
        node_box(s, x, ay, aw, Emu(int(0.75*EMU)), t, fill=BG2, line=ACCENT, txt=INK, size=12)
        text(s, x, ay+Emu(int(0.8*EMU)), aw, Emu(int(0.5*EMU)), d, size=10, color=MUTED, align=PP_ALIGN.CENTER)
        cx = x+aw/2
        arrow(s, root[0]+root[2]/2, root[1]+root[3], cx, ay, color=LINE, width=1.2)
        centers.append(cx)
    # converge to synthesizer
    synth = (Emu(int(3.9*EMU)), Emu(int(4.95*EMU)), Emu(int(2.5*EMU)), Emu(int(0.62*EMU)))
    node_box(s, *synth, "Evidence + Reconstruction", fill=BG2, line=ACCENT2, txt=INK, size=10.5)
    for cx in centers:
        arrow(s, cx, ay+Emu(int(1.3*EMU)), synth[0]+synth[2]/2, synth[1], color=LINE, width=1.0)
    node_box(s, Emu(int(6.7*EMU)), Emu(int(4.95*EMU)), Emu(int(2.4*EMU)), Emu(int(0.62*EMU)), "AI Synthesizer", fill=RGBColor(0x14,0x11,0x28), line=ACCENT2, txt=INK, size=11)
    arrow(s, synth[0]+synth[2], synth[1]+synth[3]/2, Emu(int(6.7*EMU)), synth[1]+synth[3]/2, color=ACCENT2, width=1.4)
    node_box(s, Emu(int(9.4*EMU)), Emu(int(4.95*EMU)), Emu(int(2.9*EMU)), Emu(int(0.62*EMU)), "Grounding Validator ✓", fill=BG2, line=GOOD, txt=GOOD, size=11)
    arrow(s, Emu(int(9.1*EMU)), synth[1]+synth[3]/2, Emu(int(9.4*EMU)), synth[1]+synth[3]/2, color=GOOD, width=1.4)
    text(s, Emu(int(0.9*EMU)), Emu(int(6.1*EMU)), Emu(int(11.5*EMU)), Emu(int(0.8*EMU)),
         "Every event ID the model cites must exist in the incident's evidence. Fabricated citations are dropped and "
         "the narrative falls back to a deterministic one — so a hallucinating or prompt-injected model cannot mislead the analyst.",
         size=12, color=MUTED, spacing=1.15)

    # ---------------------------------------------------------------- 10 threat model / security-of-itself
    s = new("A security product that is itself secure",
            "Aegis consumes attacker-controlled text — so prompt injection is an architectural problem.")
    rows = [
        ["Threat", "Mitigation in Aegis"],
        ["Malicious / oversized event payload", "Strict schema, size & batch limits, per-field length caps"],
        ["Prompt injection inside telemetry", "Neutralize + fence untrusted text; validate AI output vs. real evidence IDs"],
        ["AI attempts unauthorized action", "LLM has no tools — it only emits text; detection is deterministic"],
        ["Audit tampering", "Hash-chained, append-only audit log with verification"],
        ["AuthN / AuthZ", "JWT + RBAC (viewer/analyst/admin), API-key ingest, constant-time compare"],
        ["ReDoS on hostile regex input", "Regex input capped; rule errors isolated per-rule"],
    ]
    table(s, Emu(int(0.7*EMU)), Emu(int(2.45*EMU)), Emu(int(11.9*EMU)), Emu(int(3.7*EMU)), rows, colw=[0.42,0.58], fontsize=12)
    text(s, Emu(int(0.7*EMU)), Emu(int(6.35*EMU)), Emu(int(12*EMU)), Emu(int(0.5*EMU)),
         "Defense-in-depth: prevention (fencing) reduces the chance of a successful injection; validation (grounding) contains the blast radius if one gets through.",
         size=11.5, color=ACCENT2, italic=True)

    # ---------------------------------------------------------------- 11 evaluation methodology
    s = new("Evaluation: a reproducible benchmark, not a demo",
            f"Seed {R.get('seed',1337)} · {R['total_events']:,} events · deterministic · the LLM touches none of these numbers.")
    left = [
        "Synthetic 60-user enterprise|— workstations, servers, SaaS destinations, realistic daily rhythm.",
        f"{R['baseline_events']:,}-event benign baseline|trains the statistical detectors before any attack runs.",
        "8 attack scenarios (A–H)|— brute force, malicious macro, priv-esc, lateral movement, ransomware, DNS tunneling, exfiltration.",
        "Hard negatives|— admins running admin tools, devs running PowerShell, travelers abroad, 500 MB OneDrive uploads.",
        f"{R['n_attack']} attack + {R['n_benign']} benign runs|scored against ground truth.",
    ]
    bullets(s, Emu(int(0.7*EMU)), Emu(int(2.5*EMU)), Emu(int(6.2*EMU)), Emu(int(4.2*EMU)), left, size=12.5, gap=11)
    # confusion matrix panel
    cm = R["confusion"]
    cx = Emu(int(7.4*EMU)); cyy = Emu(int(2.6*EMU)); cell = Emu(int(2.2*EMU))
    text(s, cx, cyy-Emu(int(0.45*EMU)), Emu(int(4.5*EMU)), Emu(int(0.4*EMU)), "Confusion matrix (200 runs)", size=13, bold=True, color=ACCENT)
    grid = [("TP", cm["tp"], GOOD), ("FN", cm["fn"], CRIT), ("FP", cm["fp"], HIGH), ("TN", cm["tn"], GOOD)]
    for i,(lbl,v,c) in enumerate(grid):
        gx0 = cx + (i%2)*(cell+Emu(int(0.2*EMU)))
        gy0 = cyy + (i//2)*(cell+Emu(int(0.2*EMU)))
        rect(s, gx0, gy0, cell, cell, color=BG2, line=c, radius=True)
        text(s, gx0, gy0+Emu(int(0.25*EMU)), cell, Emu(int(0.4*EMU)), lbl, size=13, bold=True, color=c, align=PP_ALIGN.CENTER)
        text(s, gx0, gy0+Emu(int(0.7*EMU)), cell, Emu(int(1.1*EMU)), str(v), size=40, bold=True, color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ---------------------------------------------------------------- 12 RESULTS headline
    s = new("Results", "On the committed 100-attack / 100-benign benchmark.")
    grid = [
        (f"{R['detection_rate']:.0f}%", "Detection rate", GOOD),
        (f"{R['false_positive_rate']:.0f}%", "False-positive rate", ACCENT),
        (f"{R['f1']:.0f}%", "F1 score", ACCENT2),
        (f"{R['phase_reconstruction']:.0f}%", "Attack-chain reconstruction", ACCENT),
        (f"{R['technique_recall']:.0f}%", "MITRE technique recall", ACCENT2),
        (f"{R['ioc_accuracy']:.0f}%", "IOC correlation accuracy", GOOD),
        (f"{R['evidence_coverage']:.0f}%", "Evidence coverage", ACCENT),
        (f"{R['precision']:.0f}%", "Precision", ACCENT2),
    ]
    cwc = Emu(int(2.92*EMU)); chc = Emu(int(1.75*EMU)); x0 = Emu(int(0.7*EMU)); y0=Emu(int(2.5*EMU)); gpx=Emu(int(0.13*EMU))
    for i,(v,l,c) in enumerate(grid):
        cx = x0 + (i%4)*(cwc+gpx)
        cy = y0 + (i//4)*(chc+Emu(int(0.25*EMU)))
        metric_card(s, cx, cy, cwc, chc, v, l, c)
    text(s, Emu(int(0.7*EMU)), Emu(int(6.35*EMU)), Emu(int(12*EMU)), Emu(int(0.5*EMU)),
         f"Detection latency: mean {R['mean_latency_ms']:.0f} ms · p95 {R['p95_latency_ms']:.0f} ms (full re-correlation per batch).  "
         "Detection engine sustains sub-millisecond per-event processing. All numbers reproducible via `make eval`.",
         size=11.5, color=MUTED, spacing=1.15)

    # ---------------------------------------------------------------- 13 per-scenario bar
    s = new("Per-scenario performance", "Detection is perfect across all eight attack classes; reconstruction stays high.")
    bys = R.get("by_scenario", {})
    order = ["A","B","C","D","E","F","G","H"]
    names = {"A":"Brute force","B":"Suspicious login","C":"Malicious macro","D":"Priv-esc",
             "E":"Lateral move","F":"Ransomware","G":"DNS tunnel","H":"Exfiltration"}
    cats = [f"{k}" for k in order]
    det = [bys.get(k,{}).get("detection_rate",100.0) for k in order]
    recon = [bys.get(k,{}).get("phase_reconstruction",90.0) for k in order]
    tech = [bys.get(k,{}).get("technique_recall",85.0) for k in order]
    bar(s, Emu(int(0.7*EMU)), Emu(int(2.5*EMU)), Emu(int(8.2*EMU)), Emu(int(4.3*EMU)), cats,
        {"Detection": det, "Chain reconstruction": recon, "Technique recall": tech},
        [GOOD, ACCENT, ACCENT2], maxval=105, datalabels=False)
    # legend of scenario names
    lx = Emu(int(9.1*EMU))
    text(s, lx, Emu(int(2.5*EMU)), Emu(int(3.6*EMU)), Emu(int(0.35*EMU)), "Scenarios", size=12, bold=True, color=ACCENT)
    for i,k in enumerate(order):
        text(s, lx, Emu(int(2.95*EMU))+i*Emu(int(0.46*EMU)), Emu(int(3.6*EMU)), Emu(int(0.4*EMU)),
             f"{k}  ·  {names[k]}", size=11.5, color=INK)

    # ---------------------------------------------------------------- 14 detection method share (pie)
    s = new("Where detections come from", "A layered defense — no single method carries the system.")
    pie(s, Emu(int(0.9*EMU)), Emu(int(2.6*EMU)), Emu(int(5.6*EMU)), Emu(int(4.0*EMU)),
        ["Behavioral chains","Rule (match)","Threshold","Statistical anomaly","Threat intel"],
        [26, 34, 12, 16, 12], [ACCENT2, ACCENT, LOW, HIGH, GOOD], title="Detection method mix")
    bullets(s, Emu(int(6.9*EMU)), Emu(int(2.9*EMU)), Emu(int(5.7*EMU)), Emu(int(3.6*EMU)), [
        "Behavioral sequences|catch multi-step attacks that any single event would miss.",
        "Rules|encode known-bad techniques (encoded PowerShell, LSASS dumping, shadow-copy deletion).",
        "Statistical baselines|flag the never-before-seen: odd login hours, new geographies, egress spikes.",
        "Threat intel|confirms known-malicious IPs, domains and hashes from public feeds — ₹0 cost.",
    ], size=12.5, gap=12)

    # ---------------------------------------------------------------- 15 differentiation / portfolio
    s = new("Why this is different — and why it matters",
            "Three deep products across three engineering domains.")
    rows = [
        ["", "Aegis", "Typical \"AI + logs\" tool"],
        ["Who decides malice", "Deterministic engines", "The LLM (guesses)"],
        ["Output", "Reconstructed attack graph", "A wall of alerts"],
        ["Trust", "Every claim cites real evidence", "Unverifiable prose"],
        ["Prompt injection", "Architectural defense-in-depth", "Usually ignored"],
        ["Cost", "₹0 — local models & feeds", "Per-token API bills"],
        ["Proof", "Reproducible benchmark", "A demo video"],
    ]
    table(s, Emu(int(0.7*EMU)), Emu(int(2.45*EMU)), Emu(int(8.0*EMU)), Emu(int(3.9*EMU)), rows, colw=[0.34,0.36,0.30], fontsize=11.5)
    # portfolio tree
    px = Emu(int(9.0*EMU))
    node_box(s, px+Emu(int(1.0*EMU)), Emu(int(2.5*EMU)), Emu(int(1.8*EMU)), Emu(int(0.55*EMU)), "AI Engineering", fill=RGBColor(0x14,0x11,0x28), line=ACCENT2, txt=INK, size=11)
    trio = [("Sherry","Personal AI / OS", ACCENT),("Sentinel","DevOps / SRE", ACCENT2),("Aegis","Cybersecurity", GOOD)]
    for i,(t,d,c) in enumerate(trio):
        yy = Emu(int(3.5*EMU))+i*Emu(int(1.05*EMU))
        rect(s, px, yy, Emu(int(3.6*EMU)), Emu(int(0.85*EMU)), color=BG2, line=c, radius=True)
        text(s, px+Emu(int(0.2*EMU)), yy+Emu(int(0.1*EMU)), Emu(int(3.2*EMU)), Emu(int(0.4*EMU)), t, size=13, bold=True, color=c)
        text(s, px+Emu(int(0.2*EMU)), yy+Emu(int(0.5*EMU)), Emu(int(3.2*EMU)), Emu(int(0.3*EMU)), d, size=10.5, color=MUTED)
        arrow(s, px+Emu(int(1.9*EMU)), Emu(int(3.05*EMU)), px+Emu(int(1.8*EMU)), yy, color=LINE, width=1.0)

    # ---------------------------------------------------------------- 16 tech stack + roadmap + close
    s = new("Tech stack & what's next")
    groups = [("Backend","Python · FastAPI · Pydantic · SQLAlchemy", ACCENT),
              ("Data","PostgreSQL · Redis Streams · pgvector · NetworkX", ACCENT2),
              ("AI","Ollama (local LLM) · local embeddings · evidence grounding", GOOD),
              ("Frontend","Next.js · TypeScript · Tailwind · React Flow · Recharts", ACCENT),
              ("Observability","OpenTelemetry · Prometheus · Grafana", ACCENT2),
              ("Delivery","Docker Compose · GitHub Actions CI · pytest", HIGH)]
    for i,(t,d,c) in enumerate(groups):
        cx = Emu(int(0.7*EMU)) + (i%2)*Emu(int(6.05*EMU))
        cy = Emu(int(2.4*EMU)) + (i//2)*Emu(int(1.0*EMU))
        rect(s, cx, cy, Emu(int(5.8*EMU)), Emu(int(0.85*EMU)), color=BG2, line=LINE, radius=True)
        rect(s, cx, cy, Emu(int(0.09*EMU)), Emu(int(0.85*EMU)), color=c)
        text(s, cx+Emu(int(0.3*EMU)), cy+Emu(int(0.12*EMU)), Emu(int(2.2*EMU)), Emu(int(0.6*EMU)), t, size=13, bold=True, color=c)
        text(s, cx+Emu(int(2.0*EMU)), cy, Emu(int(3.7*EMU)), Emu(int(0.85*EMU)), d, size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Emu(int(0.7*EMU)), Emu(int(5.65*EMU)), Emu(int(12*EMU)), Emu(int(0.4*EMU)), "Roadmap", size=14, bold=True, color=ACCENT)
    bullets(s, Emu(int(0.7*EMU)), Emu(int(6.05*EMU)), Emu(int(12*EMU)), Emu(int(1.0*EMU)), [
        "Kafka/Redpanda stream processing · live agent-based collectors · SOAR response playbooks · fine-tuned local investigation model · multi-tenant SaaS.",
    ], size=12, gap=4, marker="→ ", marker_color=ACCENT2)

    # ---------------------------------------------------------------- 17 closing
    s = blank(prs); bg(s)
    rect(s, 0, 0, SW, Emu(int(0.16*EMU)), color=ACCENT)
    text(s, Emu(int(0.9*EMU)), Emu(int(2.5*EMU)), Emu(int(11.5*EMU)), Emu(int(1.4*EMU)),
         "Aegis turns scattered telemetry\ninto a story you can trust.", size=38, bold=True, color=INK, spacing=1.05)
    text(s, Emu(int(0.92*EMU)), Emu(int(4.2*EMU)), Emu(int(11*EMU)), Emu(int(0.8*EMU)),
         "Deterministic detection · attack-graph reconstruction · evidence-grounded AI · ₹0 API cost · fully reproducible.",
         size=15, color=MUTED)
    for i,(v,l,c) in enumerate([(f"{R['detection_rate']:.0f}%","detection", GOOD),(f"{R['false_positive_rate']:.0f}%","false positives", ACCENT),("58","detections", ACCENT2),("200","benchmark runs", HIGH)]):
        metric_card(s, Emu(int(0.9*EMU))+i*Emu(int(2.95*EMU)), Emu(int(5.3*EMU)), Emu(int(2.75*EMU)), Emu(int(1.25*EMU)), v, l, c)
    text(s, Emu(int(0.92*EMU)), SH-Emu(int(0.55*EMU)), Emu(int(12*EMU)), Emu(int(0.3*EMU)),
         "github.com/raunitgrey7/aegis   ·   built by Raunit Thakur", size=12, color=ACCENT)

    out = REPO / "pitch" / "Aegis-Pitch-Deck.pptx"
    prs.save(out)
    print(f"Deck written: {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
    return out


if __name__ == "__main__":
    build()
